#!/usr/bin/env python3
"""SwapShip — Final Project presentation (12 slides, premium black + cyan theme)."""

from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

# Premium dark theme: black + cyan (+ subtle purple accent)
BG = RGBColor(6, 10, 18)
BG_CARD = RGBColor(14, 20, 32)
BG_CARD_ALT = RGBColor(20, 28, 44)
CYAN = RGBColor(0, 212, 255)
CYAN_SOFT = RGBColor(120, 230, 255)
PURPLE = RGBColor(124, 58, 237)
GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(148, 163, 184)
LINE = RGBColor(51, 65, 85)

FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 12
OUT = "/Users/abhishekkumar/Desktop/SwapShip/SwapShip_Final_Project.pptx"

# --- Edit these for your submission ---
STUDENT_NAME = "Abhishek Kumar"
REG_NO = "12300520"
COLLEGE = "Your College / University Name"
DEPARTMENT = "Computer Science & Engineering"
GUIDE = "Guide / Mentor Name"
COURSE = "MVC Programming (INT221)"
PRESENTATION_DATE = date.today().strftime("%d %B %Y")
DEPLOYED = "https://swapship.onrender.com"


def fill(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency
    shape.line.fill.background()


def stroke(shape, color=LINE, pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def set_font(p, *, name=FONT_BODY, size=14, bold=False, color=WHITE, align=None):
    p.font.name = name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align


def transition(slide, kind="fade"):
    el = slide._element
    for c in list(el):
        if c.tag.endswith("transition"):
            el.remove(c)
    t = OxmlElement("p:transition")
    t.set("spd", "med")
    node = OxmlElement("p:fade" if kind == "fade" else f"p:{kind}")
    if kind in ("push", "wipe"):
        node.set("dir", "l")
    t.append(node)
    el.append(t)


def notes(slide, text):
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def dark_bg(slide, accent=True):
    fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H), BG)
    if accent:
        glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(5.5), Inches(5.5))
        fill(glow, PURPLE, 0.82)
        glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5), Inches(4), Inches(4))
        fill(glow2, CYAN, 0.88)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.05))
    fill(bar, BG_CARD)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.05), SLIDE_W, Inches(0.05))
    fill(strip, CYAN)


def footer(slide, n):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(10), Inches(0.3))
    set_font(tb.text_frame.paragraphs[0], size=9, color=MUTED)
    tb.text_frame.paragraphs[0].text = f"SwapShip  |  {COURSE}  |  {STUDENT_NAME}  |  {REG_NO}"
    pg = slide.shapes.add_textbox(Inches(12.35), Inches(7.08), Inches(0.7), Inches(0.3))
    set_font(pg.text_frame.paragraphs[0], size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    pg.text_frame.paragraphs[0].text = f"{n} / {TOTAL}"


def slide_heading(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12), Inches(0.8))
    tf = tb.text_frame
    set_font(tf.paragraphs[0], name=FONT_TITLE, size=28, bold=True, color=WHITE)
    tf.paragraphs[0].text = title
    if subtitle:
        p2 = tf.add_paragraph()
        set_font(p2, size=12, color=CYAN_SOFT)
        p2.text = subtitle


def bullets(slide, items, x=Inches(0.75), y=Inches(1.55), w=Inches(12), size=14, color=WHITE, spacing=8):
    box = slide.shapes.add_textbox(x, y, w, Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        set_font(p, size=size, color=color)
        p.space_after = Pt(spacing)


def icon_card(slide, x, y, w, h, icon, title, desc, accent=CYAN):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill(card, BG_CARD)
    stroke(card, accent, 1.2)
    card.adjustments[0] = 0.08
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.18), Inches(0.55), Inches(0.55))
    fill(circ, accent)
    it = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.22), Inches(0.55), Inches(0.45))
    set_font(it.text_frame.paragraphs[0], size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    it.text_frame.paragraphs[0].text = icon
    tt = slide.shapes.add_textbox(x + Inches(0.85), y + Inches(0.15), w - Inches(1), Inches(0.4))
    set_font(tt.text_frame.paragraphs[0], size=13, bold=True, color=WHITE)
    tt.text_frame.paragraphs[0].text = title
    dt = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.75), w - Inches(0.35), h - Inches(0.85))
    tf = dt.text_frame
    tf.word_wrap = True
    set_font(tf.paragraphs[0], size=10, color=MUTED)
    tf.paragraphs[0].text = desc


def flow_steps(slide, steps, y=Inches(2.35), x0=Inches(0.75), total_w=Inches(11.8)):
    n = len(steps)
    gap = total_w / n
    for i, (label, sub) in enumerate(steps):
        x = x0 + i * gap
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.05), y, gap - Inches(0.12), Inches(1.35))
        fill(box, BG_CARD_ALT if i % 2 else BG_CARD)
        stroke(box, CYAN if i == 0 or i == n - 1 else LINE, 1)
        box.adjustments[0] = 0.1
        lb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12), gap - Inches(0.3), Inches(0.35))
        set_font(lb.text_frame.paragraphs[0], size=11, bold=True, color=CYAN)
        lb.text_frame.paragraphs[0].text = label
        sb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.5), gap - Inches(0.3), Inches(0.75))
        tf = sb.text_frame
        tf.word_wrap = True
        set_font(tf.paragraphs[0], size=9, color=MUTED)
        tf.paragraphs[0].text = sub
        if i < n - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + gap - Inches(0.08), y + Inches(0.55), Inches(0.22), Inches(0.18))
            fill(arr, GOLD)


def content_panel(slide, x, y, w, h, title=None):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill(panel, BG_CARD)
    stroke(panel, CYAN, 1)
    panel.adjustments[0] = 0.04
    if title:
        hd = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), w - Inches(0.3), Inches(0.38))
        set_font(hd.text_frame.paragraphs[0], size=13, bold=True, color=CYAN)
        hd.text_frame.paragraphs[0].text = title
    return y + (Inches(0.55) if title else Inches(0.12))


def database_slide(slide):
    """Two-column database slide — fits safely above footer."""
    y0 = Inches(1.38)
    lh = Inches(5.55)
    lx, lw = Inches(0.65), Inches(6.15)
    rx, rw = Inches(6.95), Inches(5.75)

    content_panel(slide, lx, y0, lw, lh, "Core Tables")
    tables = [
        ("users", "Accounts, roles, verification"),
        ("items", "Listings, geo, condition, price"),
        ("item_images", "Cloudinary image URLs per item"),
        ("exchange_requests", "Swap deals between two users"),
        ("messages", "Chat per exchange (attachments)"),
        ("shipments", "AWB, status, pickup & delivery"),
        ("shipment_events", "Tracking timeline events"),
        ("orders", "Razorpay payment stages"),
    ]
    ty = y0 + Inches(0.58)
    for name, desc in tables:
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx + Inches(0.2), ty, lw - Inches(0.4), Inches(0.58))
        fill(row, BG_CARD_ALT)
        stroke(row, LINE, 0.4)
        nt = slide.shapes.add_textbox(lx + Inches(0.32), ty + Inches(0.08), Inches(1.85), Inches(0.4))
        set_font(nt.text_frame.paragraphs[0], size=10, bold=True, color=CYAN)
        nt.text_frame.paragraphs[0].text = name
        dt = slide.shapes.add_textbox(lx + Inches(2.15), ty + Inches(0.08), lw - Inches(2.4), Inches(0.42))
        set_font(dt.text_frame.paragraphs[0], size=10, color=WHITE)
        dt.text_frame.paragraphs[0].text = desc
        ty += Inches(0.6)

    content_panel(slide, rx, y0, rw, lh, "Relationships & Design")
    rel_items = [
        "One user → many items (1:N)",
        "One item → many item_images (1:N)",
        "Exchange request links requester & owner (N:1 each)",
        "Messages belong to an exchange request (1:N)",
        "One exchange → one shipment; many shipment_events (1:N)",
        "Orders track payment linked to exchange workflow",
        "Foreign keys + indexes for query performance",
        "MySQL / PostgreSQL with Laravel Eloquent ORM",
    ]
    box = slide.shapes.add_textbox(rx + Inches(0.2), y0 + Inches(0.58), rw - Inches(0.35), lh - Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(rel_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        set_font(p, size=11, color=WHITE)
        p.space_after = Pt(7)


def project_outcomes_slide(slide):
    """Slide 10 — deployment & results (no screenshots)."""
    y0 = Inches(1.38)
    content_panel(slide, Inches(0.65), y0, Inches(5.9), Inches(5.55), "Project Outcomes")
    outcomes = [
        "Fully functional P2P exchange web app deployed on Render",
        "End-to-end flow: register → list → exchange → chat → pay → ship",
        "Email OTP verification + Google OAuth for secure onboarding",
        "Real-time chat with Pusher; images via Cloudinary CDN",
        "Razorpay split payments with webhook verification",
        "Admin dashboard for users, items, orders & delivery OTP stats",
    ]
    ob = slide.shapes.add_textbox(Inches(0.85), y0 + Inches(0.58), Inches(5.5), Inches(4.8))
    tf = ob.text_frame
    tf.word_wrap = True
    for i, item in enumerate(outcomes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {item}"
        set_font(p, size=12, color=WHITE)
        p.space_after = Pt(10)

    content_panel(slide, Inches(6.75), y0, Inches(5.95), Inches(5.55), "Deployment & Demo")
    demo_items = [
        ("Live App", DEPLOYED),
        ("GitHub", "github.com/Abhii0018/SwapShip"),
        ("Health Check", f"{DEPLOYED}/healthz"),
        ("Stack", "Laravel 13 · Docker · Render"),
    ]
    dy = y0 + Inches(0.65)
    for label, value in demo_items:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), dy, Inches(5.55), Inches(0.95))
        fill(card, BG_CARD_ALT)
        stroke(card, PURPLE, 0.8)
        card.adjustments[0] = 0.12
        lt = slide.shapes.add_textbox(Inches(7.1), dy + Inches(0.1), Inches(1.4), Inches(0.35))
        set_font(lt.text_frame.paragraphs[0], size=11, bold=True, color=CYAN)
        lt.text_frame.paragraphs[0].text = label
        vt = slide.shapes.add_textbox(Inches(7.1), dy + Inches(0.42), Inches(5.25), Inches(0.45))
        set_font(vt.text_frame.paragraphs[0], size=10, color=MUTED)
        vt.text_frame.paragraphs[0].text = value
        dy += Inches(1.08)

    flow = slide.shapes.add_textbox(Inches(6.95), dy + Inches(0.15), Inches(5.55), Inches(1.2))
    tf = flow.text_frame
    tf.word_wrap = True
    set_font(tf.paragraphs[0], size=11, bold=True, color=CYAN)
    tf.paragraphs[0].text = "Suggested Demo Flow"
    steps = "Register → Add Item → Send Request → Chat → Pay → Track Shipment → Verify OTP"
    p2 = tf.add_paragraph()
    set_font(p2, size=10, color=WHITE)
    p2.text = steps
    p2.space_before = Pt(6)


def challenge_table(slide):
    headers = ["Challenge", "Solution"]
    rows = [
        ("Real-time communication", "Pusher + Laravel Echo for live chat"),
        ("Image storage at scale", "Cloudinary CDN uploads"),
        ("Secure authentication", "Laravel Breeze, email OTP, Google OAuth"),
        ("OTP email on cloud hosting", "SendGrid HTTP API (Render-safe)"),
        ("Shipment workflow complexity", "Structured shipment module + event timeline"),
        ("Payment trust", "Razorpay escrow-style split payments"),
    ]
    y = Inches(1.75)
    tw = Inches(11.9)
    cw = [Inches(4.2), Inches(7.7)]
    x0 = Inches(0.7)
    for c, h in enumerate(headers):
        x = x0 + (cw[0] if c == 0 else cw[0])
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0 if c == 0 else x0 + cw[0], y, cw[c], Inches(0.42))
        fill(cell, CYAN)
        tb = slide.shapes.add_textbox((x0 if c == 0 else x0 + cw[0]) + Inches(0.1), y + Inches(0.07), cw[c], Inches(0.35))
        set_font(tb.text_frame.paragraphs[0], size=11, bold=True, color=BG)
        tb.text_frame.paragraphs[0].text = h
    for r, row in enumerate(rows):
        ry = y + Inches(0.45) + r * Inches(0.72)
        tint = BG_CARD if r % 2 == 0 else BG_CARD_ALT
        for c, val in enumerate(row):
            x = x0 + (0 if c == 0 else cw[0])
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ry, cw[c], Inches(0.68))
            fill(cell, tint)
            stroke(cell, LINE, 0.5)
            tb = slide.shapes.add_textbox(x + Inches(0.12), ry + Inches(0.12), cw[c] - Inches(0.2), Inches(0.5))
            set_font(tb.text_frame.paragraphs[0], size=10, color=WHITE)
            tb.text_frame.paragraphs[0].text = val


def tech_stack_grid(slide):
    stacks = [
        ("Frontend", ["Blade templates", "Tailwind CSS", "Alpine.js", "Vite"], "Responsive UI & fast builds"),
        ("Backend", ["PHP 8.3", "Laravel 13", "MVC architecture"], "Secure routing, auth, ORM"),
        ("Database", ["MySQL / PostgreSQL"], "Relational data for users, items, orders"),
        ("Integrations", ["Pusher", "Cloudinary", "SendGrid", "Razorpay", "Docker"], "Chat, media, mail, payments, deploy"),
    ]
    x0, y0 = Inches(0.7), Inches(1.65)
    w, h = Inches(2.9), Inches(2.35)
    for i, (title, techs, why) in enumerate(stacks):
        col, row = i % 2, i // 2
        x = x0 + col * (w + Inches(0.35))
        y = y0 + row * (h + Inches(0.25))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        fill(card, BG_CARD)
        stroke(card, PURPLE if i % 2 else CYAN, 1.2)
        card.adjustments[0] = 0.06
        hd = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12), w, Inches(0.35))
        set_font(hd.text_frame.paragraphs[0], size=13, bold=True, color=CYAN)
        hd.text_frame.paragraphs[0].text = title
        body = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.48), w - Inches(0.25), Inches(1.0))
        tf = body.text_frame
        tf.word_wrap = True
        for j, t in enumerate(techs):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            set_font(p, size=11, color=WHITE)
            p.text = f"• {t}"
            p.space_after = Pt(2)
        why_b = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.55), w - Inches(0.25), Inches(0.65))
        tf2 = why_b.text_frame
        tf2.word_wrap = True
        set_font(tf2.paragraphs[0], size=9, color=MUTED, bold=True)
        tf2.paragraphs[0].text = f"Why: {why}"


def feature_cards(slide):
    features = [
        ("Auth", "OTP + Google OAuth"),
        ("Listings", "Upload & explore items"),
        ("Chat", "Real-time messaging"),
        ("Ship", "Tracking & delivery OTP"),
        ("Pay", "Razorpay integration"),
        ("Admin", "Dashboard & moderation"),
    ]
    x0, y0 = Inches(0.7), Inches(1.7)
    w, h = Inches(1.95), Inches(1.15)
    gap = Inches(0.22)
    for i, (title, desc) in enumerate(features):
        col, row = i % 3, i // 3
        x = x0 + col * (w + gap)
        y = y0 + row * (h + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        fill(card, BG_CARD_ALT)
        stroke(card, CYAN, 0.8)
        card.adjustments[0] = 0.1
        t1 = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.15), w, Inches(0.35))
        set_font(t1.text_frame.paragraphs[0], size=12, bold=True, color=CYAN)
        t1.text_frame.paragraphs[0].text = title
        t2 = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.5), w - Inches(0.2), Inches(0.55))
        tf = t2.text_frame
        tf.word_wrap = True
        set_font(tf.paragraphs[0], size=9, color=MUTED)
        tf.paragraphs[0].text = desc
    extra = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(11.9), Inches(0.5))
    set_font(extra.text_frame.paragraphs[0], size=11, color=WHITE)
    extra.text_frame.paragraphs[0].text = (
        "+ Search & filters  ·  Saved searches  ·  Exchange requests  ·  Responsive dark UI  ·  Deployed on Render"
    )


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ──
    s = prs.slides.add_slide(blank)
    transition(s, "fade")
    fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H), BG)
    fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.6), SLIDE_W, Inches(1.9),), BG_CARD)
    glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0.2), Inches(5.5), Inches(5.5))
    fill(glow, PURPLE, 0.78)
    glow2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4.5), Inches(4), Inches(4))
    fill(glow2, CYAN, 0.85)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.35), SLIDE_W, Inches(0.06))
    fill(line, CYAN)
    # Logo mark
    logo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.65), Inches(1.1), Inches(1.1))
    fill(logo, CYAN)
    logo.adjustments[0] = 0.2
    lt = s.shapes.add_textbox(Inches(0.85), Inches(1.82), Inches(1.1), Inches(0.7))
    set_font(lt.text_frame.paragraphs[0], size=28, bold=True, color=BG, align=PP_ALIGN.CENTER)
    lt.text_frame.paragraphs[0].text = "S"
    tb = s.shapes.add_textbox(Inches(2.15), Inches(1.55), Inches(10.5), Inches(3.8))
    tf = tb.text_frame
    set_font(tf.paragraphs[0], name=FONT_TITLE, size=56, bold=True, color=WHITE)
    tf.paragraphs[0].text = "SwapShip"
    p = tf.add_paragraph()
    set_font(p, size=20, color=CYAN)
    p.text = "Peer-to-Peer Item Exchange Platform"
    p.space_before = Pt(8)
    for text, sz, col in [
        (f"Presented by: {STUDENT_NAME}  |  Reg. {REG_NO}", 14, MUTED),
        (COLLEGE, 13, WHITE),
        (f"Department: {DEPARTMENT}", 13, MUTED),
        (f"Guide: {GUIDE}  |  {COURSE}", 13, MUTED),
        (PRESENTATION_DATE, 12, CYAN_SOFT),
    ]:
        px = tf.add_paragraph()
        set_font(px, size=sz, color=col)
        px.text = text
        px.space_before = Pt(10)
    url = s.shapes.add_textbox(Inches(2.15), Inches(5.95), Inches(8), Inches(0.35))
    set_font(url.text_frame.paragraphs[0], size=11, color=CYAN)
    url.text_frame.paragraphs[0].text = DEPLOYED
    notes(s, "Introduce yourself, college, guide, and project in 30 seconds.")
    footer(s, 1)

    # ── Slide 2: Problem Statement ──
    s = prs.slides.add_slide(blank)
    transition(s, "push")
    dark_bg(s)
    slide_heading(s, "Problem Statement", "Real-world challenges in item exchange")
    bullets(
        s,
        [
            "People often have unused items but lack an easy platform to exchange them.",
            "Existing platforms focus on buying/selling — not peer-to-peer swapping.",
            "Users face trust issues: payment risk, fake listings, unverified strangers.",
            "Communication and shipment tracking are fragmented across multiple apps.",
            "No single system covers discovery → chat → payment → delivery proof.",
        ],
        y=Inches(1.5),
        size=13,
    )
    icon_card(s, Inches(0.75), Inches(5.15), Inches(3.6), Inches(1.55), "⇄", "Exchange Gap", "No dedicated swap-first marketplace", CYAN)
    icon_card(s, Inches(4.55), Inches(5.15), Inches(3.6), Inches(1.55), "📦", "Shipping Pain", "Tracking & handover lack transparency", PURPLE)
    icon_card(s, Inches(8.35), Inches(5.15), Inches(3.6), Inches(1.55), "💬", "Trust & Chat", "Negotiation happens outside secure channels", GOLD)
    notes(s, "Explain why swap/trust/shipping matter to everyday users.")
    footer(s, 2)

    # ── Slide 3: Project Overview ──
    s = prs.slides.add_slide(blank)
    transition(s, "fade")
    dark_bg(s)
    slide_heading(s, "Project Overview", "What is SwapShip?")
    bullets(
        s,
        [
            "SwapShip is a web-based item exchange platform for students and general users.",
            "Users upload items, explore listings, send exchange requests, chat, pay, and track shipments.",
            "Built as a full-stack Laravel MVC academic project with production deployment.",
        ],
        y=Inches(1.45),
        w=Inches(5.9),
        size=12,
    )
    mod = s.shapes.add_textbox(Inches(0.75), Inches(3.05), Inches(5.9), Inches(2.2))
    tf = mod.text_frame
    set_font(tf.paragraphs[0], size=12, bold=True, color=CYAN)
    tf.paragraphs[0].text = "Core Modules"
    modules = [
        "User Authentication (OTP + Google)",
        "Item Management & Explore",
        "Exchange Requests",
        "Real-Time Messaging",
        "Shipment Tracking",
        "Payment System (Razorpay)",
        "Admin Dashboard",
    ]
    for m in modules:
        p = tf.add_paragraph()
        set_font(p, size=11, color=MUTED)
        p.text = f"• {m}"
        p.space_after = Pt(2)
    # Workflow diagram (right panel)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.45), Inches(5.65), Inches(5.15))
    fill(panel, BG_CARD)
    stroke(panel, CYAN, 1)
    panel.adjustments[0] = 0.04
    pt = s.shapes.add_textbox(Inches(7.2), Inches(1.55), Inches(5.2), Inches(0.4))
    set_font(pt.text_frame.paragraphs[0], size=12, bold=True, color=CYAN)
    pt.text_frame.paragraphs[0].text = "Workflow"
    flow_steps(
        s,
        [
            ("User", "Register"),
            ("Upload", "List item"),
            ("Request", "Exchange"),
            ("Ship", "Track"),
            ("Done", "Complete"),
        ],
        y=Inches(2.15),
        x0=Inches(7.15),
        total_w=Inches(5.35),
    )
    notes(s, "One-liner: SwapShip = trusted swap marketplace end-to-end.")
    footer(s, 3)

    # ── Slide 4: Objectives ──
    s = prs.slides.add_slide(blank)
    transition(s, "wipe")
    dark_bg(s)
    slide_heading(s, "Objectives of the Project")
    bullets(
        s,
        [
            "Build a secure item exchange platform with verified users.",
            "Enable smooth shipment management with status events & delivery OTP.",
            "Provide real-time communication between buyers and sellers.",
            "Create a responsive, modern UI for better user experience.",
            "Maintain secure authentication, authorization, and admin oversight.",
            "Reduce complexity of traditional multi-app exchange workflows.",
        ],
        y=Inches(1.55),
        size=14,
    )
    notes(s, "Read objectives clearly — maps to evaluation criteria.")
    footer(s, 4)

    # ── Slide 5: Technology Stack ──
    s = prs.slides.add_slide(blank)
    transition(s, "fade")
    dark_bg(s)
    slide_heading(s, "Technology Stack", "Tools chosen and why")
    tech_stack_grid(s)
    notes(s, "Mention Laravel MVC, Tailwind responsiveness, Pusher real-time, Cloudinary images.")
    footer(s, 5)

    # ── Slide 6: System Architecture / Workflow ──
    s = prs.slides.add_slide(blank)
    transition(s, "push")
    dark_bg(s)
    slide_heading(s, "System Architecture & Workflow", "End-to-end transaction flow")
    flow_steps(
        s,
        [
            ("1", "Register / Login"),
            ("2", "Upload Item"),
            ("3", "Exchange Req."),
            ("4", "Chat"),
            ("5", "Payment"),
            ("6", "Shipment"),
            ("7", "Complete"),
        ],
        y=Inches(1.55),
    )
    arch = s.shapes.add_textbox(Inches(0.75), Inches(3.15), Inches(11.8), Inches(2.0))
    tf = arch.text_frame
    set_font(tf.paragraphs[0], size=12, bold=True, color=CYAN)
    tf.paragraphs[0].text = "MVC Architecture"
    layers = [
        "View: Blade + Tailwind + Alpine.js  →  Controller: Item, Exchange, Message, Shipment, Payment, Admin",
        "Model: User, Item, ExchangeRequest, Message, Shipment, Order  →  Services: Shipping, OTP Mail, Payments",
    ]
    for i, line in enumerate(layers):
        p = tf.add_paragraph()
        set_font(p, size=11, color=MUTED)
        p.text = line
        p.space_after = Pt(8)
    notes(s, "Walk the flowchart left-to-right — this is a key evaluator slide.")
    footer(s, 6)

    # ── Slide 7: Database Design ──
    s = prs.slides.add_slide(blank)
    transition(s, "fade")
    dark_bg(s, accent=False)
    slide_heading(s, "Database Design", "Schema & relationships")
    database_slide(s)
    notes(s, "Explain normalization and foreign keys briefly.")
    footer(s, 7)

    # ── Slide 8: Key Features ──
    s = prs.slides.add_slide(blank)
    transition(s, "wipe")
    dark_bg(s)
    slide_heading(s, "Key Features", "Major capabilities of SwapShip")
    feature_cards(s)
    notes(s, "Optionally show 1 live screenshot from swapship.onrender.com.")
    footer(s, 8)

    # ── Slide 9: Challenges & Solutions ──
    s = prs.slides.add_slide(blank)
    transition(s, "fade")
    dark_bg(s)
    slide_heading(s, "Challenges Faced & Solutions", "Problem-solving during development")
    challenge_table(s)
    notes(s, "Interviewers value this slide — show you debugged real issues.")
    footer(s, 9)

    # ── Slide 10: Project Outcomes & Deployment ──
    s = prs.slides.add_slide(blank)
    transition(s, "push")
    dark_bg(s, accent=False)
    slide_heading(s, "Project Outcomes & Deployment", "Results, live links & demo flow")
    project_outcomes_slide(s)
    notes(s, "Share live URL, walk through demo flow, mention Render + Docker deployment.")
    footer(s, 10)

    # ── Slide 11: Future Enhancements ──
    s = prs.slides.add_slide(blank)
    transition(s, "fade")
    dark_bg(s)
    slide_heading(s, "Future Enhancements", "Scalability & vision")
    bullets(
        s,
        [
            "AI-based item recommendations and smart matching.",
            "Native mobile application (Android / iOS).",
            "Advanced shipment tracking with courier API integrations.",
            "Rating & review system for user reputation.",
            "Video call / richer media support in chat.",
            "Multi-language support for wider adoption.",
        ],
        y=Inches(1.55),
        size=14,
    )
    notes(s, "Shows forward thinking without undermining current scope.")
    footer(s, 11)

    # ── Slide 12: Conclusion & Q/A ──
    s = prs.slides.add_slide(blank)
    transition(s, "fade")
    fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H), BG)
    fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.05), SLIDE_W, Inches(0.06)), CYAN)
    # Corner accents only — keep center clear for readable text
    glow_l = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4.5), Inches(4.5), Inches(4.5))
    fill(glow_l, CYAN, 0.82)
    glow_r = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5), Inches(4.5), Inches(4.5))
    fill(glow_r, PURPLE, 0.82)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(1.45), Inches(10.65), Inches(5.35))
    fill(card, BG_CARD)
    stroke(card, CYAN, 1.5)
    card.adjustments[0] = 0.03
    tb = s.shapes.add_textbox(Inches(1.65), Inches(1.75), Inches(10.05), Inches(4.85))
    tf = tb.text_frame
    tf.word_wrap = True
    set_font(tf.paragraphs[0], name=FONT_TITLE, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].text = "Conclusion"
    p2 = tf.add_paragraph()
    set_font(p2, size=14, color=CYAN_SOFT, align=PP_ALIGN.CENTER)
    p2.text = (
        "SwapShip provides a modern, secure platform for item exchange with integrated "
        "shipment management, real-time communication, and scalable Laravel MVC architecture."
    )
    p2.space_before = Pt(16)
    p3 = tf.add_paragraph()
    set_font(p3, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    p3.text = (
        "Achieved: verified users · escrow payments · live tracking · admin oversight · cloud deployment"
    )
    p3.space_before = Pt(14)
    p4 = tf.add_paragraph()
    set_font(p4, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    p4.text = "Laravel 13 · Tailwind · MySQL · Pusher · Cloudinary · SendGrid · Razorpay · Docker · Render"
    p4.space_before = Pt(12)
    p5 = tf.add_paragraph()
    set_font(p5, size=26, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    p5.text = "Thank You"
    p5.space_before = Pt(22)
    p6 = tf.add_paragraph()
    set_font(p6, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    p6.text = "Questions?"
    p6.space_before = Pt(8)
    p7 = tf.add_paragraph()
    set_font(p7, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    p7.text = f"{STUDENT_NAME}  ·  {REG_NO}  ·  {COURSE}"
    p7.space_before = Pt(14)
    notes(s, "Summarize impact, invite questions confidently.")
    footer(s, 12)

    prs.save(OUT)
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
